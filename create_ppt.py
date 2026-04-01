import os
import re
import base64
import requests
import pandas as pd
import matplotlib.pyplot as plt
from dotenv import load_dotenv
from openai import OpenAI
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# ==================== 0. 基础准备 ====================
os.makedirs("data", exist_ok=True)

# ==================== 1. 加载环境变量和初始化客户端 ====================
load_dotenv()

OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
MINIMAX_API_KEY = os.getenv("MINIMAX_API_KEY")

if not OPENROUTER_API_KEY:
    raise ValueError("未找到 OpenRouter API Key，请在 .env 中设置 OPENROUTER_API_KEY")

client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=OPENROUTER_API_KEY,
)

MODEL = "stepfun/step-3.5-flash:free"

# ==================== 2. 读取和处理 CSV 数据 ====================
file_path = "data/sales_data.csv"
if not os.path.exists(file_path):
    raise FileNotFoundError(f"未找到数据文件：{file_path}")

df = pd.read_csv(file_path)

required_col = "日期"
if required_col not in df.columns:
    raise ValueError(f"CSV 中缺少必要列：{required_col}")

# 解析日期（格式：日/月/年）
df["日期"] = pd.to_datetime(df["日期"], format="%d/%m/%Y", dayfirst=True, errors="raise")

# 提取年份和季度
df["year"] = df["日期"].dt.year
df["quarter"] = df["日期"].dt.quarter

# 获取图书列名（除了日期、year、quarter以外的列）
book_columns = [col for col in df.columns if col not in ["日期", "year", "quarter"]]
if not book_columns:
    raise ValueError("未检测到产品列，请检查 sales_data.csv 的表头")

# 转长表
df_long = pd.melt(
    df,
    id_vars=["year", "quarter"],
    value_vars=book_columns,
    var_name="product",
    value_name="sales",
)

# 确保 sales 为数值
df_long["sales"] = pd.to_numeric(df_long["sales"], errors="coerce").fillna(0)

# 按年、季度、图书汇总
quarterly_sales = (
    df_long.groupby(["year", "quarter", "product"], as_index=False)["sales"]
    .sum()
    .sort_values(["year", "quarter", "product"])
)

quarterly_sales["year_quarter"] = (
    quarterly_sales["year"].astype(str) + " Q" + quarterly_sales["quarter"].astype(str)
)

print("=== 处理后的季度销售数据（前5行）===")
print(quarterly_sales.head())

# ==================== 3. 生成销售额折线图 ====================
plot_img_path = "data/季度销售额折线图.png"

plt.figure(figsize=(12, 6))

for book in book_columns:
    book_data = quarterly_sales[quarterly_sales["product"] == book]
    if not book_data.empty:
        plt.plot(
            book_data["year_quarter"],
            book_data["sales"],
            marker="o",
            label=book,
        )

plt.title("2022-2025年各图书季度销售额")
plt.xlabel("季度")
plt.ylabel("销售额")
plt.xticks(rotation=45)
plt.legend()
plt.grid(True, linestyle="--", alpha=0.6)
plt.tight_layout()
plt.savefig(plot_img_path, dpi=300)
plt.close()

print(f"图表已保存为 {plot_img_path}")

# ==================== 4. 准备数据摘要 ====================
summary_text = f"""
以下是图书的季度销售数据（2022-2025年）：
- 图书：{", ".join(book_columns)}
- 季度销售额明细（年份、季度、产品、销售额）：
{quarterly_sales[["year", "quarter", "product", "sales"]].to_string(index=False)}
""".strip()

# ==================== 5. 工具函数：调用 LLM ====================
def ask_llm(messages, model=MODEL):
    """
    调用 OpenRouter / OpenAI SDK 兼容接口
    """
    response = client.chat.completions.create(
        model=model,
        messages=messages,
        temperature=0.3,
    )
    return response.choices[0].message.content.strip()


def clean_title(text: str) -> str:
    """
    清洗标题，避免模型返回引号、编号、Markdown
    """
    if not text:
        return "销售趋势概览"

    text = text.strip()
    text = re.sub(r"^[\"'“”‘’]+|[\"'“”‘’]+$", "", text)
    text = re.sub(r"^\s*\d+[\.、]\s*", "", text)
    text = text.replace("#", "").strip()

    # 过长则截断
    if len(text) > 20:
        text = text[:20]

    return text or "销售趋势概览"


def parse_insights(text: str):
    """
    把模型返回解析成项目符号列表
    """
    if not text:
        return ["销售整体呈现阶段性变化", "重点品类在部分季度增长明显"]

    lines = [line.strip() for line in text.splitlines() if line.strip()]
    cleaned = []

    for line in lines:
        line = re.sub(r"^\s*\d+[\.、]\s*", "", line)
        line = line.replace("•", "").replace("-", "").strip()
        if line:
            cleaned.append(line)

    if not cleaned:
        cleaned = [text.strip()]

    # 最多保留 3 条
    return cleaned[:3]


# ==================== 6. 多轮对话演示：鲜花价格（独立会话） ====================
price_messages = []

user_input_1 = "你能够帮我计算鲜花的价格，3束玫瑰每束12元，再加2束百合每束15元，一共多少钱？"
price_messages.append({"role": "user", "content": user_input_1})
assistant_reply_1 = ask_llm(price_messages)
print("助手（计算价格）:", assistant_reply_1)
price_messages.append({"role": "assistant", "content": assistant_reply_1})

user_input_2 = "我把每个花束定价为进价基础上加价20%,进价80元时,我的售价是多少。"
price_messages.append({"role": "user", "content": user_input_2})
assistant_reply_2 = ask_llm(price_messages)
print("助手（加价）:", assistant_reply_2)
price_messages.append({"role": "assistant", "content": assistant_reply_2})

# ==================== 7. 生成洞察和标题（独立会话，避免上下文串场） ====================
analysis_messages = [
    {
        "role": "system",
        "content": (
            "你是一个商业分析助手。"
            "你的任务是根据销售数据生成适合 PPT 展示的简洁洞察与标题。"
            "输出必须简洁、专业、中文。"
        ),
    }
]

insight_prompt = f"""
请根据以下销售数据摘要，给出两个最重要的洞察。
要求：
1. 每条一句话
2. 每句尽量控制在 18~24 个字
3. 用数字编号
4. 不要解释过程，只输出结果

{summary_text}
""".strip()

analysis_messages.append({"role": "user", "content": insight_prompt})
insights_raw = ask_llm(analysis_messages)
print("模型生成的洞察：", insights_raw)
analysis_messages.append({"role": "assistant", "content": insights_raw})

title_prompt = """
根据你刚才总结的洞察，为这页销售分析幻灯片生成一个非常简短的标题。
要求：
1. 不超过 12 个字
2. 适合商务 PPT
3. 只输出标题本身
""".strip()

analysis_messages.append({"role": "user", "content": title_prompt})
title_raw = ask_llm(analysis_messages)
print("模型生成的标题：", title_raw)
analysis_messages.append({"role": "assistant", "content": title_raw})

title = clean_title(title_raw)
insight_list = parse_insights(insights_raw)

# ==================== 8. 调用 MiniMax 生成公司图片 ====================
company_summary = "我们是网络鲜花批发商，但是我们董事长也写IT图书！"
cover_img_path = "data/花语秘境咖哥.png"

def save_placeholder_image(path):
    img = Image.new("RGB", (1024, 1024), color=(73, 109, 137))
    img.save(path)

def generate_image_with_minimax(output_path: str, summary: str):
    if not MINIMAX_API_KEY:
        print("未设置 MINIMAX_API_KEY，使用纯色占位图片")
        save_placeholder_image(output_path)
        return

    url = "https://api.minimaxi.com/v1/image_generation"
    headers = {
        "Authorization": f"Bearer {MINIMAX_API_KEY}",
        "Content-Type": "application/json",
    }

    prompt = (
        f"根据这个公司概述：{summary}。"
        "请创建一张适合季度销售规划会议封面的高质量商业图片，"
        "画面体现成长、前进、品牌感、鲜花行业气质与现代商业氛围，"
        "整体真实、高级、适合PPT封面。"
    )

    payload = {
        "model": "image-01",
        "prompt": prompt,
        "aspect_ratio": "1:1",
        "response_format": "base64",   # 也可改成 "url"
        "n": 1,
        "prompt_optimizer": True,
    }

    try:
        resp = requests.post(url, headers=headers, json=payload, timeout=60)
        resp.raise_for_status()

        result = resp.json()
        print("MiniMax 返回：", result)

        base_resp = result.get("base_resp", {})
        if base_resp.get("status_code") != 0:
            raise RuntimeError(f"MiniMax 业务失败: {base_resp}")

        data = result.get("data", {})
        image_base64_list = data.get("image_base64")
        image_urls = data.get("image_urls")

        if image_base64_list and len(image_base64_list) > 0:
            img_bytes = base64.b64decode(image_base64_list[0])
            with open(output_path, "wb") as f:
                f.write(img_bytes)
            print("已通过 MiniMax API 生成图片（base64）并保存:", output_path)
            return

        if image_urls and len(image_urls) > 0:
            img_response = requests.get(image_urls[0], timeout=60)
            img_response.raise_for_status()
            with open(output_path, "wb") as f:
                f.write(img_response.content)
            print("已通过 MiniMax API 生成图片（URL）并保存:", output_path)
            return

        raise RuntimeError(f"MiniMax 返回格式异常，未找到 image_base64 或 image_urls: {result}")

    except Exception as e:
        print(f"调用 MiniMax 图片生成 API 失败: {e}")
        print("使用纯色占位图片代替")
        save_placeholder_image(output_path)

generate_image_with_minimax(cover_img_path, company_summary)

# ==================== 9. 创建 PPTX 幻灯片 ====================
def add_text_paragraph(text_frame, text, font_size=18, bold=False, color=(255, 255, 255), level=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    return p


def create_slide1(prs, image_path, title_text, subtitle_text):
    """
    第一张幻灯片：公司图片 + 标题 + 副标题
    """
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # 背景黑色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 左侧图片
    img_left = 0
    img_top = 0
    img_width = int(slide_width * 0.6)
    img_height = slide_height
    slide.shapes.add_picture(image_path, img_left, img_top, width=img_width, height=img_height)

    # 标题
    title_left = int(slide_width * 0.62)
    title_top = Inches(2.0)
    title_width = int(slide_width * 0.33)
    title_height = Inches(1.2)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.clear()

    p = title_frame.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # 副标题
    subtitle_left = int(slide_width * 0.62)
    subtitle_top = Inches(3.1)
    subtitle_width = int(slide_width * 0.33)
    subtitle_height = Inches(0.8)

    subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()

    p2 = subtitle_frame.paragraphs[0]
    p2.text = subtitle_text
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(220, 220, 220)
    p2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


def create_slide2(prs, image_path, title_text, insight_items):
    """
    第二张幻灯片：折线图 + 标题 + 洞察项目符号
    """
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 顶部标题
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), slide_width - Inches(0.6), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()

    title_p = title_frame.paragraphs[0]
    title_p.text = title_text
    title_p.font.bold = True
    title_p.font.size = Pt(26)
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # 左侧图表
    chart_left = Inches(0.3)
    chart_top = Inches(1.3)
    chart_width = int(slide_width * 0.60)
    chart_height = slide_height - Inches(2.0)
    slide.shapes.add_picture(image_path, chart_left, chart_top, width=chart_width, height=chart_height)

    # 右侧洞察框
    insight_left = int(slide_width * 0.66)
    insight_top = Inches(1.4)
    insight_width = int(slide_width * 0.28)
    insight_height = Inches(4.5)

    insights_box = slide.shapes.add_textbox(insight_left, insight_top, insight_width, insight_height)
    insights_frame = insights_box.text_frame
    insights_frame.clear()

    header = insights_frame.paragraphs[0]
    header.text = "关键见解"
    header.font.bold = True
    header.font.size = Pt(22)
    header.font.color.rgb = RGBColor(0, 180, 140)

    for item in insight_items:
        p = insights_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

# ==================== 10. 生成 PPT ====================
prs = Presentation()

# 可选：设置 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

create_slide1(prs, cover_img_path, "花语秘境", "2025年销售大会")
create_slide2(prs, plot_img_path, title, insight_list)

pptx_path = "data/咖哥花语秘境.pptx"
prs.save(pptx_path)

print(f"PPTX 已生成: {pptx_path}")